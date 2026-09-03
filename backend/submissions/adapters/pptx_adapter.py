from __future__ import annotations

from io import BytesIO

from pptx import Presentation

from submissions.adapters.base import BaseSubmissionAdapter, ExtractedDocument
from submissions.text_sanitize import sanitize_json, sanitize_text


class PptxAdapter(BaseSubmissionAdapter):
    file_type = "pptx"

    def extract(self, data: bytes, filename: str) -> ExtractedDocument:
        prs = Presentation(BytesIO(data))
        slides = []
        for idx, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            slides.append({"slide": idx + 1, "texts": texts})
        full_text = "\n\n".join("\n".join(s["texts"]) for s in slides)
        return ExtractedDocument(
            text=sanitize_text(full_text),
            structure=sanitize_json({"slides": slides, "slide_count": len(slides)}),
            source_ref=filename,
        )
